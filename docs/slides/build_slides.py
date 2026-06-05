#!/usr/bin/env python3
"""
마크다운 -> PPTX 생성기 (python-pptx).

비교/기술 의사결정 덱을 위한 '작은 마크다운 부분집합'만 지원한다.
(의도적으로 단순 — 우리가 작성하는 마크다운을 이 규칙 안에서 쓰면 슬라이드가 안정적으로 나온다.)

지원 문법
  # 제목                 -> 타이틀 슬라이드 (바로 다음 줄들은 부제)
  ## 슬라이드 제목        -> 새 콘텐츠 슬라이드
  ### 굵은 리드 줄        -> 슬라이드 안 굵은 소제목
  - 불릿 / 2칸 들여쓰기   -> 불릿(2칸=하위, 4칸=하위하위)
  | a | b | + | --- |     -> 표 (GFM)
  ```                     -> 코드 블록(고정폭, 연회색 배경)
  코드 ...
  ```
  > 메모                  -> 작은 회색 노트
  인라인: **굵게**, `코드`

사용법
  python build_slides.py <input.md> [output.pptx]
"""
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# ---- 색/레이아웃 상수 -------------------------------------------------------
PRIMARY = RGBColor(0x1F, 0x38, 0x64)   # 진한 네이비 (제목)
ACCENT = RGBColor(0x2E, 0x6B, 0xE6)    # 파랑 (강조/코드)
GRAY = RGBColor(0x66, 0x66, 0x66)      # 노트
DARK = RGBColor(0x22, 0x22, 0x22)      # 본문
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG = RGBColor(0xF2, 0xF4, 0xF7)   # 코드 블록 배경

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_L = Inches(0.7)
CONTENT_W = Inches(13.333 - 1.4)
BLANK = 6  # 기본 템플릿의 빈 레이아웃 인덱스

INLINE = re.compile(r"(\*\*.+?\*\*|`.+?`)")


def strip_inline(text):
    """표 셀 등 런 분리가 필요 없는 곳: 마크업 기호만 제거."""
    return text.replace("**", "").replace("`", "")


def add_runs(paragraph, text, size, color=DARK):
    """**굵게**, `코드` 인라인을 런 단위로 렌더."""
    for tok in INLINE.split(text):
        if not tok:
            continue
        run = paragraph.add_run()
        if tok.startswith("**") and tok.endswith("**"):
            run.text = tok[2:-2]
            run.font.bold = True
            run.font.color.rgb = color
        elif tok.startswith("`") and tok.endswith("`"):
            run.text = tok[1:-1]
            run.font.name = "Consolas"
            run.font.color.rgb = ACCENT
        else:
            run.text = tok
            run.font.color.rgb = color
        run.font.size = size


def set_indent(paragraph, level):
    """블랭크 텍스트박스에서도 들여쓰기가 보이도록 marL을 직접 지정."""
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("marL", str(int(Inches(0.3) * level)))
    pPr.set("indent", "0")


# ---- 마크다운 파싱 ----------------------------------------------------------
def parse(md):
    lines = md.splitlines()
    title_slide = None       # (title, subtitle)
    slides = []              # [{'title':..., 'blocks':[(kind, payload), ...]}]
    current = None
    i, n = 0, len(lines)

    while i < n:
        line = lines[i]
        s = line.strip()

        if s.startswith("# ") and title_slide is None and current is None:
            title = s[2:].strip()
            i += 1
            sub = []
            while i < n and lines[i].strip() and not lines[i].strip().startswith("#"):
                sub.append(lines[i].strip())
                i += 1
            title_slide = (title, " ".join(sub))
            continue

        if s.startswith("## "):
            current = {"title": s[3:].strip(), "blocks": []}
            slides.append(current)
            i += 1
            continue

        if current is None:
            i += 1
            continue

        if s.startswith("```"):
            code = []
            i += 1
            while i < n and not lines[i].strip().startswith("```"):
                code.append(lines[i])
                i += 1
            i += 1  # 닫는 ``` 소비
            current["blocks"].append(("code", code))
            continue

        if s.startswith("### "):
            current["blocks"].append(("lead", s[4:].strip()))
            i += 1
            continue

        if s.startswith("> "):
            current["blocks"].append(("note", s[2:].strip()))
            i += 1
            continue

        if s.startswith("|"):
            raw = []
            while i < n and lines[i].strip().startswith("|"):
                raw.append(lines[i].strip())
                i += 1
            rows = []
            for r in raw:
                cells = [c.strip() for c in r.strip().strip("|").split("|")]
                if all(set(c) <= set("-: ") and "-" in c for c in cells):
                    continue  # 구분선 행
                rows.append(cells)
            current["blocks"].append(("table", rows))
            continue

        if s.startswith("- ") or s.startswith("* "):
            bullets = []
            while i < n and (lines[i].lstrip().startswith("- ") or lines[i].lstrip().startswith("* ")):
                raw = lines[i]
                indent = (len(raw) - len(raw.lstrip(" "))) // 2
                bullets.append((indent, raw.lstrip()[2:].strip()))
                i += 1
            current["blocks"].append(("bullets", bullets))
            continue

        if s == "":
            i += 1
            continue

        current["blocks"].append(("paragraph", s))
        i += 1

    return title_slide, slides


# ---- 렌더링 -----------------------------------------------------------------
def render_title(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), SLIDE_W, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    box = slide.shapes.add_textbox(MARGIN_L, Inches(2.7), CONTENT_W, Inches(1.6))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    add_runs(p, title, Pt(40), PRIMARY)
    for r in p.runs:
        r.font.bold = True
    if subtitle:
        sp = tf.add_paragraph()
        sp.space_before = Pt(10)
        add_runs(sp, subtitle, Pt(18), GRAY)


def render_title_bar(slide, title):
    box = slide.shapes.add_textbox(MARGIN_L, Inches(0.4), CONTENT_W, Inches(0.85))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    add_runs(p, title, Pt(27), PRIMARY)
    for r in p.runs:
        r.font.bold = True
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_L, Inches(1.22), Inches(2.2), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def render_block(slide, block, top):
    kind, payload = block

    if kind == "bullets":
        h = Inches(0.36) * len(payload) + Inches(0.1)
        box = slide.shapes.add_textbox(MARGIN_L, top, CONTENT_W, h)
        tf = box.text_frame
        tf.word_wrap = True
        for idx, (level, text) in enumerate(payload):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.space_after = Pt(5)
            set_indent(p, level)
            marker = "•  " if level == 0 else "–  "
            lead = paragraph_prefix(p, marker, level)
            add_runs(p, text, Pt(18) if level == 0 else Pt(16))
        return top + h + Inches(0.18)

    if kind in ("paragraph", "lead"):
        box = slide.shapes.add_textbox(MARGIN_L, top, CONTENT_W, Inches(0.5))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        add_runs(p, payload, Pt(18))
        if kind == "lead":
            for r in p.runs:
                r.font.bold = True
                r.font.color.rgb = PRIMARY
        return top + Inches(0.5)

    if kind == "note":
        box = slide.shapes.add_textbox(MARGIN_L, top, CONTENT_W, Inches(0.45))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        add_runs(p, payload, Pt(13), GRAY)
        for r in p.runs:
            r.font.italic = True
        return top + Inches(0.45)

    if kind == "code":
        h = Inches(0.26) * len(payload) + Inches(0.2)
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN_L, top, CONTENT_W, h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = CODE_BG
        rect.line.color.rgb = RGBColor(0xD5, 0xDB, 0xE5)
        tf = rect.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.18)
        tf.margin_top = Inches(0.08)
        for idx, ln in enumerate(payload):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = ln if ln else " "
            run.font.name = "Consolas"
            run.font.size = Pt(12)
            run.font.color.rgb = DARK
        return top + h + Inches(0.18)

    if kind == "table":
        rows = payload
        nrows = len(rows)
        ncols = max(len(r) for r in rows)
        h = Inches(0.38) * nrows
        gframe = slide.shapes.add_table(nrows, ncols, MARGIN_L, top, CONTENT_W, h)
        table = gframe.table
        col_w = int(CONTENT_W / ncols)
        for ci in range(ncols):
            table.columns[ci].width = col_w
        for ri, row in enumerate(rows):
            for ci in range(ncols):
                cell = table.cell(ri, ci)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.margin_top = Inches(0.03)
                cell.margin_bottom = Inches(0.03)
                cell.text = strip_inline(row[ci]) if ci < len(row) else ""
                para = cell.text_frame.paragraphs[0]
                for run in para.runs:
                    run.font.size = Pt(13)
                    if ri == 0:
                        run.font.bold = True
                        run.font.color.rgb = WHITE
                    else:
                        run.font.color.rgb = DARK
                if ri == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = PRIMARY
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = WHITE if ri % 2 else RGBColor(0xF6, 0xF8, 0xFB)
        return top + h + Inches(0.25)

    return top


def paragraph_prefix(paragraph, marker, level):
    """불릿 마커를 별도 런으로 추가(색은 강조색)."""
    run = paragraph.add_run()
    run.text = marker
    run.font.size = Pt(18) if level == 0 else Pt(16)
    run.font.color.rgb = ACCENT
    run.font.bold = True
    return run


def build(md_path, out_path):
    md = Path(md_path).read_text(encoding="utf-8")
    title_slide, slides = parse(md)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    if title_slide:
        render_title(prs, *title_slide)
    for s in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
        render_title_bar(slide, s["title"])
        top = Inches(1.5)
        for block in s["blocks"]:
            top = render_block(slide, block, top)

    prs.save(out_path)
    total = (1 if title_slide else 0) + len(slides)
    print(f"OK  {out_path}  ({total} slides)")


def main():
    if len(sys.argv) < 2:
        print("usage: python build_slides.py <input.md> [output.pptx]")
        sys.exit(1)
    md_path = Path(sys.argv[1])
    out_path = Path(sys.argv[2]) if len(sys.argv) > 2 else md_path.with_suffix(".pptx")
    build(md_path, out_path)


if __name__ == "__main__":
    main()
